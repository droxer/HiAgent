"""Document generation tools for creating PDF, DOCX, XLSX, and PPTX files in the sandbox."""

from __future__ import annotations

import base64
import re
import shlex
from pathlib import PurePosixPath
from typing import Any

from agent.tools.base import (
    ExecutionContext,
    SandboxTool,
    ToolDefinition,
    ToolResult,
)

_SCRIPT_PATH = "/tmp/_doc_gen_script.py"
_CONTENT_PATH = "/tmp/_doc_content.md"
_TEMPLATE_DIR = "/opt/doc_templates"
_VALID_STYLES = frozenset({"default", "report", "article", "minimal"})
_VALID_PAGE_SIZES = frozenset({"A4", "letter", "legal"})
_VALID_ORIENTATIONS = frozenset({"portrait", "landscape"})
_OUTPUT_PATH_RE = re.compile(r"^/workspace/[\w./\-]+$")


def _validate_output_path(output_path: str) -> str | None:
    """Return an error message if output_path is unsafe, or None if OK."""
    resolved = PurePosixPath(output_path)
    if ".." in resolved.parts:
        return "output_path must not contain '..'"
    if not _OUTPUT_PATH_RE.match(output_path):
        return "output_path must be within /workspace/ and contain only safe characters"
    return None


async def _write_script_to_sandbox(session: Any, script: str) -> bool:
    """Write the generation script into the sandbox, returning True on success.

    Tries ``write_file`` first, then falls back to base64-encoded ``exec``
    if the file does not appear in the sandbox.
    """
    try:
        await session.write_file(_SCRIPT_PATH, script)
    except OSError:
        pass  # fall through to verification

    check = await session.exec(f"test -f {shlex.quote(_SCRIPT_PATH)}")
    if check.success:
        return True

    # Fallback: write the script via base64-decoded exec (avoids copy_in)
    encoded = base64.b64encode(script.encode("utf-8")).decode("ascii")
    write_result = await session.exec(
        f"echo {shlex.quote(encoded)} | base64 -d > {shlex.quote(_SCRIPT_PATH)}"
    )
    return write_result.success


async def _run_doc_script(
    session: Any, script: str, output_path: str
) -> ToolResult:
    """Write and execute a document generation script, returning the result."""
    if not await _write_script_to_sandbox(session, script):
        return ToolResult.fail(
            "Failed to write generation script to sandbox"
        )

    result = await session.exec(
        f"python3 {shlex.quote(_SCRIPT_PATH)}", timeout=120
    )

    if result.exit_code != 0:
        error = result.stderr or result.stdout or "Unknown error"
        return ToolResult.fail(f"Document generation failed: {error}")

    # Verify file was created — shell-quote the path to prevent injection
    quoted = shlex.quote(output_path)
    check = await session.exec(f"test -f {quoted} && stat -c %s {quoted}")
    if check.exit_code != 0:
        return ToolResult.fail(f"Output file not created at {output_path}")

    size = check.stdout.strip()
    return ToolResult.ok(
        f"Document created at {output_path} ({size} bytes)",
        metadata={
            "artifact_paths": [output_path],
            "path": output_path,
            "size": int(size) if size.isdigit() else 0,
        },
    )


def _validate_pdf_params(
    style: str, page_size: str, orientation: str
) -> str | None:
    """Return an error message if any parameter is invalid, or None if all OK."""
    if style not in _VALID_STYLES:
        return f"Invalid style '{style}'. Must be one of: {', '.join(sorted(_VALID_STYLES))}"
    if page_size not in _VALID_PAGE_SIZES:
        return f"Invalid page_size '{page_size}'. Must be one of: {', '.join(sorted(_VALID_PAGE_SIZES))}"
    if orientation not in _VALID_ORIENTATIONS:
        return f"Invalid orientation '{orientation}'. Must be one of: {', '.join(sorted(_VALID_ORIENTATIONS))}"
    return None


# The generation script reads content and custom CSS from files, and uses
# pre-baked HTML/CSS templates. This avoids all string escaping issues.
#
# Two-level brace escaping:
#   - str.format() consumes one level: {{ → {  and }} → }
#   - The generated Python code uses f-strings with single braces
#   - CSS template placeholders like {{PAGE_SIZE}} need four braces: {{{{PAGE_SIZE}}}}
#     → str.format() produces {{PAGE_SIZE}} → CSS .replace() substitutes the value
_PDF_SCRIPT_TEMPLATE = """\
import sys
import pathlib

CONTENT_PATH = "{content_path}"
TEMPLATE_DIR = pathlib.Path("{template_dir}")
OUTPUT_PATH = "{output_path}"
STYLE = "{style}"
PAGE_SIZE = "{page_size}"
ORIENTATION = "{orientation}"
TITLE = "{title}"
CUSTOM_CSS_PATH = "{custom_css_path}"

# --- Read content and custom CSS from files (no escaping needed) ---
content_md = pathlib.Path(CONTENT_PATH).read_text(encoding="utf-8")
custom_css = pathlib.Path(CUSTOM_CSS_PATH).read_text(encoding="utf-8")

# --- Convert markdown to HTML ---
import markdown
html_body = markdown.markdown(
    content_md,
    extensions=["tables", "fenced_code", "codehilite", "toc", "attr_list"],
    extension_configs={{
        "codehilite": {{"css_class": "codehilite", "guess_lang": True}},
    }},
)

# --- Load template and CSS ---
base_html = (TEMPLATE_DIR / "base.html").read_text(encoding="utf-8")
css_text = (TEMPLATE_DIR / "styles" / f"{{STYLE}}.css").read_text(encoding="utf-8")

# --- Apply page size and orientation to CSS ---
css_text = css_text.replace("{{{{PAGE_SIZE}}}}", PAGE_SIZE)
css_text = css_text.replace("{{{{ORIENTATION}}}}", ORIENTATION)
css_text = css_text.replace("{{{{TITLE}}}}", TITLE)

# --- Build title block ---
title_block = ""
if TITLE:
    title_block = f'<h1 class="doc-title">{{TITLE}}</h1>'

# --- Build custom CSS block ---
custom_css_block = ""
if custom_css.strip():
    custom_css_block = f"<style>{{custom_css}}</style>"

# --- Assemble final HTML ---
final_html = base_html
final_html = final_html.replace("{{{{TITLE}}}}", TITLE)
final_html = final_html.replace("{{{{CSS}}}}", css_text)
final_html = final_html.replace("{{{{CUSTOM_CSS}}}}", custom_css_block)
final_html = final_html.replace("{{{{TITLE_BLOCK}}}}", title_block)
final_html = final_html.replace("{{{{CONTENT}}}}", html_body)

# --- Render PDF ---
from weasyprint import HTML
HTML(string=final_html).write_pdf(OUTPUT_PATH)
print(f"PDF created: {{OUTPUT_PATH}}")
"""


class DocCreatePDF(SandboxTool):
    """Create a styled PDF document from markdown content using WeasyPrint."""

    def definition(self) -> ToolDefinition:
        return ToolDefinition(
            name="document_create_pdf",
            description=(
                "Create a styled PDF document from markdown content. Supports "
                "headings, lists, tables, code blocks with syntax highlighting, "
                "bold/italic, blockquotes, and images. Choose from built-in "
                "themes (default, report, article, minimal) or provide custom CSS."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "content": {
                        "type": "string",
                        "description": (
                            "Content for the PDF in markdown format. Supports "
                            "headings, lists, tables, code blocks, bold/italic, "
                            "blockquotes, and images."
                        ),
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output file path in the sandbox.",
                        "default": "/workspace/output.pdf",
                    },
                    "title": {
                        "type": "string",
                        "description": "Optional document title displayed at the top.",
                        "default": "",
                    },
                    "style": {
                        "type": "string",
                        "description": (
                            "Visual theme. 'default': clean sans-serif; "
                            "'report': formal with page numbers and header; "
                            "'article': wide margins, large text, drop caps; "
                            "'minimal': no decoration, max density."
                        ),
                        "enum": ["default", "report", "article", "minimal"],
                        "default": "default",
                    },
                    "custom_css": {
                        "type": "string",
                        "description": (
                            "Optional CSS rules appended after the theme stylesheet. "
                            "Use for one-off styling overrides."
                        ),
                        "default": "",
                    },
                    "page_size": {
                        "type": "string",
                        "description": "Page size for the PDF.",
                        "enum": ["A4", "letter", "legal"],
                        "default": "A4",
                    },
                    "orientation": {
                        "type": "string",
                        "description": "Page orientation.",
                        "enum": ["portrait", "landscape"],
                        "default": "portrait",
                    },
                },
                "required": ["content"],
            },
            execution_context=ExecutionContext.SANDBOX,
            tags=("document", "pdf", "sandbox"),
        )

    async def execute(self, session: Any, **kwargs: Any) -> ToolResult:
        content: str = kwargs.get("content", "")
        output_path: str = kwargs.get("output_path", "/workspace/output.pdf")
        title: str = kwargs.get("title", "")
        style: str = kwargs.get("style", "default")
        custom_css: str = kwargs.get("custom_css", "")
        page_size: str = kwargs.get("page_size", "A4")
        orientation: str = kwargs.get("orientation", "portrait")

        if not content.strip():
            return ToolResult.fail("Content must not be empty")

        path_error = _validate_output_path(output_path)
        if path_error is not None:
            return ToolResult.fail(path_error)

        validation_error = _validate_pdf_params(style, page_size, orientation)
        if validation_error is not None:
            return ToolResult.fail(validation_error)

        # Write content as a separate file — avoids all string escaping issues
        await session.write_file(_CONTENT_PATH, content)

        # Write custom CSS as a separate file too — avoids triple-quote escaping
        custom_css_path = "/tmp/_doc_custom.css"
        await session.write_file(custom_css_path, custom_css)

        # Title: escape for double-quoted Python string literal, strip newlines
        safe_title = (
            title.replace("\\", "\\\\")
            .replace('"', '\\"')
            .replace("\n", " ")
            .replace("\r", "")
        )

        script = _PDF_SCRIPT_TEMPLATE.format(
            content_path=_CONTENT_PATH,
            template_dir=_TEMPLATE_DIR,
            output_path=output_path,
            style=style,
            page_size=page_size,
            orientation=orientation,
            title=safe_title,
            custom_css_path=custom_css_path,
        )

        return await _run_doc_script(session, script, output_path)


class DocCreateDOCX(SandboxTool):
    """Create a DOCX document."""

    def definition(self) -> ToolDefinition:
        return ToolDefinition(
            name="document_create_docx",
            description=(
                "Create a Word DOCX document. Provide content as text with "
                "paragraphs separated by blank lines."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "content": {
                        "type": "string",
                        "description": "Text content for the document.",
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output file path in the sandbox.",
                        "default": "/workspace/output.docx",
                    },
                    "title": {
                        "type": "string",
                        "description": "Optional document title/heading.",
                        "default": "",
                    },
                },
                "required": ["content"],
            },
            execution_context=ExecutionContext.SANDBOX,
            tags=("document", "docx", "sandbox"),
        )

    async def execute(self, session: Any, **kwargs: Any) -> ToolResult:
        content: str = kwargs.get("content", "")
        output_path: str = kwargs.get("output_path", "/workspace/output.docx")
        title: str = kwargs.get("title", "")

        if not content.strip():
            return ToolResult.fail("Content must not be empty")

        path_error = _validate_output_path(output_path)
        if path_error is not None:
            return ToolResult.fail(path_error)

        # Write content and title as separate files
        content_file = "/tmp/_docx_content.txt"
        title_file = "/tmp/_docx_title.txt"
        await session.write_file(content_file, content)
        await session.write_file(title_file, title)

        script = f'''\
import subprocess
import sys
import pathlib

try:
    from docx import Document
    from docx.shared import Pt
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-docx", "-q"])
    from docx import Document
    from docx.shared import Pt

content = pathlib.Path("{content_file}").read_text(encoding="utf-8")
title = pathlib.Path("{title_file}").read_text(encoding="utf-8")

doc = Document()
if title:
    doc.add_heading(title, level=0)

for para in content.split("\\n\\n"):
    para = para.strip()
    if para:
        p = doc.add_paragraph(para)
        p.style.font.size = Pt(11)

doc.save("{output_path}")
print(f"DOCX created: {output_path}")
'''
        return await _run_doc_script(session, script, output_path)


class DocCreateXLSX(SandboxTool):
    """Create an Excel XLSX spreadsheet."""

    def definition(self) -> ToolDefinition:
        return ToolDefinition(
            name="document_create_xlsx",
            description=(
                "Create an Excel XLSX spreadsheet. Provide data as a JSON array of arrays "
                "(rows), or as CSV-formatted text. First row is treated as headers."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "data": {
                        "type": "string",
                        "description": (
                            "Data as JSON array of arrays, e.g. "
                            '\'[["Name","Age"],["Alice",30]]\', '
                            "or as CSV text."
                        ),
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output file path in the sandbox.",
                        "default": "/workspace/output.xlsx",
                    },
                    "sheet_name": {
                        "type": "string",
                        "description": "Name of the worksheet.",
                        "default": "Sheet1",
                    },
                },
                "required": ["data"],
            },
            execution_context=ExecutionContext.SANDBOX,
            tags=("document", "xlsx", "sandbox"),
        )

    async def execute(self, session: Any, **kwargs: Any) -> ToolResult:
        data: str = kwargs.get("data", "")
        output_path: str = kwargs.get("output_path", "/workspace/output.xlsx")
        sheet_name: str = kwargs.get("sheet_name", "Sheet1")

        if not data.strip():
            return ToolResult.fail("Data must not be empty")

        path_error = _validate_output_path(output_path)
        if path_error is not None:
            return ToolResult.fail(path_error)

        # Write data and sheet name as separate files
        data_file = "/tmp/_xlsx_data.txt"
        sheet_file = "/tmp/_xlsx_sheet.txt"
        await session.write_file(data_file, data)
        await session.write_file(sheet_file, sheet_name)

        script = f'''\
import subprocess
import sys
import json
import csv
import io
import pathlib

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "openpyxl", "-q"])
    from openpyxl import Workbook
    from openpyxl.styles import Font

raw_data = pathlib.Path("{data_file}").read_text(encoding="utf-8")
sheet_name = pathlib.Path("{sheet_file}").read_text(encoding="utf-8")

# Try JSON first, then CSV
try:
    rows = json.loads(raw_data)
except json.JSONDecodeError:
    reader = csv.reader(io.StringIO(raw_data))
    rows = list(reader)

wb = Workbook()
ws = wb.active
ws.title = sheet_name

for row_idx, row in enumerate(rows, 1):
    for col_idx, value in enumerate(row, 1):
        cell = ws.cell(row=row_idx, column=col_idx, value=value)
        if row_idx == 1:
            cell.font = Font(bold=True)

wb.save("{output_path}")
print(f"XLSX created: {output_path} ({{len(rows)}} rows)")
'''
        return await _run_doc_script(session, script, output_path)


class DocCreatePPTX(SandboxTool):
    """Create a PowerPoint PPTX presentation."""

    def definition(self) -> ToolDefinition:
        return ToolDefinition(
            name="document_create_pptx",
            description=(
                "Create a PowerPoint PPTX presentation. Provide slides as a JSON array "
                "of objects with 'title' and 'content' fields, or as text with slides "
                "separated by '---'."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "slides": {
                        "type": "string",
                        "description": (
                            "Slides as JSON array, e.g. "
                            '\'[{"title":"Intro","content":"Hello world"}]\', '
                            "or as text with slides separated by '---'."
                        ),
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Output file path in the sandbox.",
                        "default": "/workspace/output.pptx",
                    },
                },
                "required": ["slides"],
            },
            execution_context=ExecutionContext.SANDBOX,
            tags=("document", "pptx", "sandbox"),
        )

    async def execute(self, session: Any, **kwargs: Any) -> ToolResult:
        slides_input: str = kwargs.get("slides", "")
        output_path: str = kwargs.get("output_path", "/workspace/output.pptx")

        if not slides_input.strip():
            return ToolResult.fail("Slides data must not be empty")

        path_error = _validate_output_path(output_path)
        if path_error is not None:
            return ToolResult.fail(path_error)

        # Write slides data as a separate file
        slides_file = "/tmp/_pptx_slides.txt"
        await session.write_file(slides_file, slides_input)

        script = f'''\
import subprocess
import sys
import json
import pathlib

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt

raw = pathlib.Path("{slides_file}").read_text(encoding="utf-8")

# Parse slides
try:
    slides_data = json.loads(raw)
except json.JSONDecodeError:
    # Parse text format: slides separated by ---
    parts = raw.split("---")
    slides_data = []
    for part in parts:
        lines = part.strip().split("\\n", 1)
        title = lines[0].strip() if lines else "Slide"
        content = lines[1].strip() if len(lines) > 1 else ""
        slides_data.append({{"title": title, "content": content}})

prs = Presentation()

for slide_info in slides_data:
    title = slide_info.get("title", "")
    content = slide_info.get("content", "")

    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if content:
        slide.placeholders[1].text = content

prs.save("{output_path}")
print(f"PPTX created: {output_path} ({{len(slides_data)}} slides)")
'''
        return await _run_doc_script(session, script, output_path)
